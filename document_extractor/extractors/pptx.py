"""
PowerPoint text extraction using python-pptx.
"""

from pathlib import Path
from typing import Union
from pptx import Presentation
from pptx.exc import PackageNotFoundError

from ..models import DocumentResult


class PPTXExtractor:
    """Extracts text from PowerPoint files using python-pptx."""
    
    @staticmethod
    def extract(file_path: Union[str, Path]) -> DocumentResult:
        """
        Extract text from a PowerPoint file.
        
        Args:
            file_path: Path to the PowerPoint file (string or Path object)
            
        Returns:
            DocumentResult: Extraction result containing the text and metadata
        """
        file_path = Path(file_path)
        
        try:
            if not file_path.exists():
                raise FileNotFoundError(f"File not found: {file_path}")
                
            if file_path.suffix.lower() != '.pptx':
                raise ValueError(f"Not a PPTX file: {file_path}")
            
            text = []
            prs = Presentation(str(file_path))
            
            # Extract text from each slide's shapes
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(shape.text.strip())
            
            return DocumentResult.from_path(
                file_path,
                content="\n".join(text)
            )
            
        except (FileNotFoundError, ValueError) as e:
            # Pass through common errors with their messages
            return DocumentResult.from_path(file_path, error=str(e))
        except PackageNotFoundError:
            return DocumentResult.from_path(
                file_path,
                error="Invalid or corrupted PPTX file"
            )
        except Exception as e:
            return DocumentResult.from_path(
                file_path,
                error=f"Unexpected error during PPTX extraction: {e}"
            )

    @staticmethod
    def extract_text(file_path: Union[str, Path]) -> str:
        """
        Legacy method for backward compatibility.
        
        Args:
            file_path: Path to the PowerPoint file (string or Path object)
            
        Returns:
            str: Extracted text from the presentation
            
        Raises:
            FileNotFoundError: If the file doesn't exist
            ValueError: If the file is not a PPTX
        """
        result = PPTXExtractor.extract(file_path)
        if result.error:
            raise ValueError(result.error)
        return result.content